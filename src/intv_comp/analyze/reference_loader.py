"""
追加資料（参考資料）の読み込みモジュール。

法制審議会の議事録等の追加資料をファイルシステムから読み込み、
LLMに提供するための機能を提供する。
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import List

from intv_comp.logger import logger

# Optional imports for document processing
try:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    PdfReadError = Exception  # type: ignore[misc, assignment]  # Fallback for exception handling

try:
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    DocxPackageNotFoundError = Exception  # type: ignore[misc, assignment]  # Fallback for exception handling

try:
    from openpyxl import load_workbook
    from openpyxl.utils.exceptions import InvalidFileException
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False
    InvalidFileException = Exception  # Fallback for exception handling

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    PptxPackageNotFoundError = Exception  # type: ignore[misc, assignment]  # Fallback for exception handling

try:
    from PIL import Image
    import pytesseract
    from pytesseract.pytesseract import TesseractNotFoundError, TesseractError
    # Verify tesseract is actually available
    pytesseract.get_tesseract_version()
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    TesseractNotFoundError = Exception  # Fallback for exception handling
    TesseractError = Exception  # Fallback for exception handling
except Exception:
    # Catch Tesseract-specific exceptions and OS errors
    # (TesseractNotFoundError, TesseractError, OSError)
    HAS_OCR = False
    TesseractNotFoundError = Exception  # Fallback for exception handling
    TesseractError = Exception  # Fallback for exception handling


def _extract_text_from_pdf(file_path: Path) -> str:
    """PDFファイルからテキストを抽出する。"""
    if not HAS_PDF:
        logger.warning(f"PDF処理ライブラリが利用できません: {file_path.name}")
        return ""

    try:
        reader = PdfReader(file_path)
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"[ページ {page_num}]\n{text}")
        return "\n\n".join(text_parts)
    except PdfReadError as exc:
        # PDF が壊れている・形式がおかしい等
        logger.warning("PDF読み込みエラー(不正なPDF): %s - %s", file_path.name, exc)
    except OSError as exc:
        # ファイルが存在しない、権限がない、I/O エラーなど
        logger.warning("PDFファイルアクセスエラー: %s - %s", file_path.name, exc)

    # エラーが発生した場合は空文字列を返す
    return ""


def _extract_text_from_docx(file_path: Path) -> str:
    """Wordファイルからテキストを抽出する。"""
    if not HAS_DOCX:
        logger.warning(f"Word処理ライブラリが利用できません: {file_path.name}")
        return ""

    try:
        doc = Document(str(file_path))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    except DocxPackageNotFoundError as exc:
        # Wordファイルが壊れている / Word形式ではない場合
        logger.warning("Word読み込みエラー(不正なWordファイル): %s - %s", file_path.name, exc)
    except OSError as exc:
        # ファイルが存在しない、権限がない、I/O エラーなど
        logger.warning("Wordファイルアクセスエラー: %s - %s", file_path.name, exc)

    # エラーが発生した場合は空文字列を返す
    return ""


def _extract_text_from_xlsx(file_path: Path) -> str:
    """Excelファイルからテキストを抽出する。"""
    if not HAS_XLSX:
        logger.warning(f"Excel処理ライブラリが利用できません: {file_path.name}")
        return ""

    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                text_parts.append(f"[シート: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(text_parts)
    except InvalidFileException as exc:
        # xlsx ではない / 壊れているなど
        logger.warning("Excel読み込みエラー(不正なExcelファイル): %s - %s", file_path.name, exc)
    except OSError as exc:
        # ファイルアクセス系エラー
        logger.warning("Excelファイルアクセスエラー: %s - %s", file_path.name, exc)

    # エラーが発生した場合は空文字列を返す
    return ""


def _extract_text_from_pptx(file_path: Path) -> str:
    """PowerPointファイルからテキストを抽出する。"""
    if not HAS_PPTX:
        logger.warning(f"PowerPoint処理ライブラリが利用できません: {file_path.name}")
        return ""

    try:
        prs = Presentation(str(file_path))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                text_parts.append(f"[スライド {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except PptxPackageNotFoundError as exc:
        # pptx が壊れている / pptx ではない場合
        logger.warning("PowerPoint読み込みエラー(不正なファイル): %s - %s", file_path.name, exc)
    except OSError as exc:
        # I/O、権限などのファイルアクセスエラー
        logger.warning("PowerPointファイルアクセスエラー: %s - %s", file_path.name, exc)

    # エラーが発生した場合は空文字列を返す
    return ""


def _extract_text_from_image(file_path: Path) -> str:
    """画像ファイルからOCRでテキストを抽出する。"""
    if not HAS_OCR:
        logger.warning(f"OCR処理ライブラリが利用できません: {file_path.name}")
        return ""

    try:
        with Image.open(str(file_path)) as image:
            text = pytesseract.image_to_string(image, lang='jpn+eng')
            return str(text).strip()
    except TesseractError as exc:
        # OCR 処理中のエラー
        logger.warning("画像OCR処理エラー: %s - %s", file_path.name, exc)
    except OSError as exc:
        # ファイルアクセス系エラー
        logger.warning("画像ファイルアクセスエラー: %s - %s", file_path.name, exc)

    # エラーが発生した場合は空文字列を返す
    return ""


def load_reference_materials(references_dir: Path) -> str:
    """
    指定されたディレクトリから追加資料を読み込み、結合したテキストを返す。

    Args:
        references_dir: 追加資料が格納されているディレクトリのパス

    Returns:
        読み込んだ全ての資料を結合したテキスト。資料がない場合は空文字列。
    """
    if not references_dir.exists():
        logger.info(f"追加資料のディレクトリが存在しません: {references_dir}")
        return ""

    if not references_dir.is_dir():
        logger.warning(f"指定されたパスはディレクトリではありません: {references_dir}")
        return ""

    # 最大ファイルサイズを環境変数から取得（デフォルト: 30MB）
    try:
        max_file_size = int(os.getenv("MAX_REFERENCE_FILE_SIZE", "31457280"))
    except ValueError:
        logger.warning("MAX_REFERENCE_FILE_SIZE is not a valid integer, using default 30MB")
        max_file_size = 31457280
    # サポートされる拡張子と対応する抽出関数
    file_processors = {
        ".txt": lambda p: p.read_text(encoding="utf-8"),
        ".md": lambda p: p.read_text(encoding="utf-8"),
        ".pdf": _extract_text_from_pdf,
        ".docx": _extract_text_from_docx,
        ".xlsx": _extract_text_from_xlsx,
        ".pptx": _extract_text_from_pptx,
        ".jpg": _extract_text_from_image,
        ".jpeg": _extract_text_from_image,
        ".png": _extract_text_from_image,
    }

    # ファイルを収集
    reference_files: List[Path] = []
    for ext in file_processors:
        reference_files.extend(references_dir.rglob(f"*{ext}"))

    if not reference_files:
        logger.info(f"追加資料が見つかりませんでした: {references_dir}")
        return ""

    # ファイルを読み込んで結合
    materials: List[str] = []
    for file_path in sorted(reference_files):
        try:
            # ファイルサイズをチェック
            file_size = file_path.stat().st_size
            if file_size > max_file_size:
                logger.warning(
                    f"ファイルサイズが大きすぎます（{file_size / 1024 / 1024:.1f}MB > "
                    f"{max_file_size / 1024 / 1024:.1f}MB）: {file_path.name}"
                )
                continue

            ext = file_path.suffix.lower()
            processor = file_processors.get(ext)

            if processor is None:
                logger.warning(f"未対応のファイル形式です: {file_path.name}")
                continue

            content = processor(file_path)
            if content and content.strip():
                materials.append(f"# {file_path.name}\n\n{content}")
                logger.info(f"追加資料を読み込みました: {file_path.name}")
            else:
                logger.warning(f"ファイルからテキストを抽出できませんでした: {file_path.name}")
        except (OSError, PermissionError, UnicodeDecodeError) as exc:
            logger.warning(f"資料の読み込みに失敗しました: {file_path.name} - {exc}")
            continue

    if not materials:
        logger.info("読み込み可能な追加資料がありませんでした")
        return ""

    combined_text = "\n\n---\n\n".join(materials)
    logger.info(f"{len(materials)}件の追加資料を読み込みました")
    return combined_text
